import os
import json
import uuid
import io
from typing import List
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from app.core.config import settings

# Try importing Gemini API and document parsers
try:
    import google.generativeai as genai
    HAS_GEMINI = True
except ImportError:
    HAS_GEMINI = False

try:
    from pypdf import PdfReader
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

BASE_DIR = os.path.dirname(os.path.dirname(os.path.dirname(__file__)))
DB_PATH = os.path.join(BASE_DIR, "vector_store_db.json")

class LocalVectorStore:
    def __init__(self):
        self.documents = []
        self.chunks = []
        self.load_db()

    def load_db(self):
        if os.path.exists(DB_PATH):
            try:
                with open(DB_PATH, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    self.documents = data.get("documents", [])
                    self.chunks = data.get("chunks", [])
            except Exception as e:
                self.documents = []
                self.chunks = []
        else:
            self.documents = []
            self.chunks = []

    def save_db(self):
        try:
            with open(DB_PATH, "w", encoding="utf-8") as f:
                json.dump({
                    "documents": self.documents,
                    "chunks": self.chunks
                }, f, ensure_ascii=False, indent=2)
        except Exception:
            pass

    def add_document(self, filename: str, file_type: str, size_bytes: int, text: str):
        doc_id = str(uuid.uuid4())
        chunk_size = 750
        overlap = 150
        chunks_extracted = []
        
        i = 0
        while i < len(text):
            chunk_text = text[i:i + chunk_size].strip()
            if chunk_text:
                chunks_extracted.append({
                    "id": str(uuid.uuid4()),
                    "doc_id": doc_id,
                    "filename": filename,
                    "text": chunk_text
                })
            i += (chunk_size - overlap)
            
        if not chunks_extracted and text.strip():
            chunks_extracted.append({
                "id": str(uuid.uuid4()),
                "doc_id": doc_id,
                "filename": filename,
                "text": text.strip()
            })

        self.documents.append({
            "id": doc_id,
            "filename": filename,
            "file_type": file_type,
            "size_bytes": size_bytes,
            "total_chunks": len(chunks_extracted)
        })
        self.chunks.extend(chunks_extracted)
        self.save_db()
        return doc_id, len(chunks_extracted)

    def delete_document(self, doc_id: str):
        self.documents = [d for d in self.documents if d["id"] != doc_id]
        self.chunks = [c for c in self.chunks if c["doc_id"] != doc_id]
        self.save_db()

    def search_similar(self, query: str, top_k: int = 3) -> List[dict]:
        if not self.chunks:
            return []
            
        chunk_texts = [c["text"] for c in self.chunks]
        
        try:
            vectorizer = TfidfVectorizer(stop_words='english')
            tfidf_matrix = vectorizer.fit_transform(chunk_texts)
            query_vec = vectorizer.transform([query])
            
            similarities = cosine_similarity(query_vec, tfidf_matrix).flatten()
            top_indices = np.argsort(similarities)[::-1][:top_k]
            
            results = []
            for idx in top_indices:
                score = float(similarities[idx])
                if score >= 0.02 or len(chunk_texts) < 5:
                    results.append({
                        "chunk": self.chunks[idx],
                        "score": score
                    })
            return results
        except Exception:
            results = []
            for chunk in self.chunks:
                score = 0.0
                words = query.lower().split()
                matches = sum(1 for w in words if w in chunk["text"].lower())
                if matches > 0:
                    score = matches / len(words)
                    results.append({"chunk": chunk, "score": score})
            results.sort(key=lambda x: x["score"], reverse=True)
            return results[:top_k]

vector_store = LocalVectorStore()

def extract_text_from_bytes(file_bytes: bytes, filename: str) -> str:
    ext = os.path.splitext(filename.lower())[1]
    
    if ext == ".pdf":
        if HAS_PYPDF:
            try:
                reader = PdfReader(io.BytesIO(file_bytes))
                text = ""
                for page in reader.pages:
                    content = page.extract_text()
                    if content:
                        text += content + "\n"
                return text
            except Exception as e:
                return f"[Error parsing PDF: {e}]"
        return "[Error: pypdf is not installed]"
            
    elif ext == ".docx":
        if HAS_DOCX:
            try:
                doc = docx.Document(io.BytesIO(file_bytes))
                return "\n".join([p.text for p in doc.paragraphs if p.text])
            except Exception as e:
                return f"[Error parsing DOCX: {e}]"
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                xml_content = z.read('word/document.xml')
                root = ET.fromstring(xml_content)
                paragraphs = [elem.text for elem in root.iter() if elem.tag.endswith('t') and elem.text]
                return " ".join(paragraphs)
        except Exception as e:
            return f"[DOCX parser error fallback: {e}]"

    elif ext == ".pptx":
        if HAS_PPTX:
            try:
                prs = Presentation(io.BytesIO(file_bytes))
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text_runs.append(shape.text)
                return "\n".join(text_runs)
            except Exception as e:
                return f"[Error parsing PPTX: {e}]"
        return "[Error: python-pptx is not installed]"
            
    try:
        return file_bytes.decode('utf-8', errors='ignore')
    except Exception as e:
        return f"[Error decoding text file: {e}]"

def get_gemini_model():
    if HAS_GEMINI and settings.GEMINI_API_KEY:
        try:
            genai.configure(api_key=settings.GEMINI_API_KEY)
            return genai.GenerativeModel('gemini-1.5-flash')
        except Exception:
            return None
    return None
